from __future__ import annotations

from dataclasses import dataclass
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "docs" / "presentations"
OUTPUT_PATH = OUTPUT_DIR / "framework-to-code-mapping-cn.pptx"
GENERATION_MANIFEST_PATH = ROOT / "projects" / "knowledge_base_basic" / "generated" / "generation_manifest.json"
GOVERNANCE_MANIFEST_PATH = ROOT / "projects" / "knowledge_base_basic" / "generated" / "governance_manifest.json"


@dataclass(frozen=True)
class Theme:
    width: Any
    height: Any
    bg: RGBColor
    paper: RGBColor
    ink: RGBColor
    muted: RGBColor
    accent: RGBColor
    accent_alt: RGBColor
    accent_dark: RGBColor
    ok: RGBColor
    warn: RGBColor
    gold: RGBColor
    title_font: str
    body_font: str
    mono_font: str


THEME = Theme(
    width=Inches(13.333),
    height=Inches(7.5),
    bg=RGBColor(247, 244, 238),
    paper=RGBColor(255, 255, 255),
    ink=RGBColor(33, 38, 45),
    muted=RGBColor(97, 103, 113),
    accent=RGBColor(232, 91, 64),
    accent_alt=RGBColor(40, 116, 118),
    accent_dark=RGBColor(41, 59, 65),
    ok=RGBColor(55, 132, 108),
    warn=RGBColor(197, 126, 40),
    gold=RGBColor(228, 175, 73),
    title_font="Aptos Display",
    body_font="Aptos",
    mono_font="Cascadia Code",
)


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def set_background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    color: RGBColor,
    *,
    rounded: bool = True,
) -> Any:
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_outline_rect(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    line_color: RGBColor,
    fill_color: RGBColor | None = None,
) -> Any:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or THEME.paper
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.2)
    return shape


def add_text(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    font_size: int,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name or THEME.body_font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or THEME.ink
    return textbox


def add_paragraphs(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    lines: Iterable[str],
    *,
    font_size: int = 14,
    color: RGBColor | None = None,
    font_name: str | None = None,
    bullet: bool = False,
) -> Any:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        if bullet:
            paragraph.level = 0
            paragraph.bullet = True
        run = paragraph.runs[0]
        run.font.name = font_name or THEME.body_font
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or THEME.ink
        paragraph.space_after = Pt(4)
    return textbox


def add_section_title(slide: Any, eyebrow: str, title: str, subtitle: str) -> None:
    add_text(slide, Inches(0.55), Inches(0.35), Inches(2.5), Inches(0.25), eyebrow, font_size=13, color=THEME.accent, bold=True)
    add_text(
        slide,
        Inches(0.55),
        Inches(0.65),
        Inches(10.8),
        Inches(0.5),
        title,
        font_size=25,
        bold=True,
        font_name=THEME.title_font,
    )
    add_text(slide, Inches(0.55), Inches(1.08), Inches(12.0), Inches(0.4), subtitle, font_size=12, color=THEME.muted)


def add_footer(slide: Any, text: str) -> None:
    add_text(slide, Inches(0.6), Inches(7.02), Inches(12.0), Inches(0.22), text, font_size=10, color=THEME.muted)


def add_connector(slide: Any, x1: Any, y1: Any, x2: Any, y2: Any, *, color: RGBColor | None = None) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color or THEME.accent_dark
    connector.line.width = Pt(1.8)


def add_arrow_connector(slide: Any, x1: Any, y1: Any, x2: Any, y2: Any, *, color: RGBColor | None = None) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color or THEME.accent_dark
    connector.line.width = Pt(2.0)


def add_card(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    title: str,
    lines: Iterable[str],
    accent: RGBColor,
) -> None:
    add_outline_rect(slide, left, top, width, height, accent)
    add_rect(slide, left, top, width, Inches(0.14), accent, rounded=False)
    add_text(slide, left + Inches(0.14), top + Inches(0.18), width - Inches(0.28), Inches(0.25), title, font_size=15, bold=True)
    add_paragraphs(
        slide,
        left + Inches(0.14),
        top + Inches(0.5),
        width - Inches(0.28),
        height - Inches(0.6),
        lines,
        font_size=12,
        color=THEME.muted,
    )


def _governed_object(payload: dict[str, Any], object_id: str) -> dict[str, Any]:
    for item in payload.get("structural_objects", []):
        if isinstance(item, dict) and item.get("object_id") == object_id:
            return item
    raise ValueError(f"missing structural object: {object_id}")


def build_cover(prs: Any, generation_manifest: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)

    add_rect(slide, Inches(0), Inches(0), Inches(4.95), Inches(7.5), THEME.accent_dark, rounded=False)
    add_text(slide, Inches(0.6), Inches(0.72), Inches(3.8), Inches(0.3), "Shelf / Framework Mapping", font_size=16, color=THEME.gold, bold=True)
    add_text(
        slide,
        Inches(0.6),
        Inches(1.1),
        Inches(3.9),
        Inches(1.1),
        "从 Framework 到 Code\n再从 Code 反查回上游",
        font_size=29,
        color=THEME.paper,
        bold=True,
        font_name=THEME.title_font,
    )
    add_paragraphs(
        slide,
        Inches(0.6),
        Inches(2.55),
        Inches(3.8),
        Inches(1.4),
        [
            "这份 PPT 对应当前 knowledge_base_basic 主链。",
            "目标不是讲概念，而是让人顺着文件追到一一对应。",
            "主线是：Framework -> Product Spec -> Implementation Config -> Code -> Evidence。",
        ],
        font_size=15,
        color=RGBColor(231, 234, 237),
    )

    add_card(
        slide,
        Inches(5.35),
        Inches(0.75),
        Inches(3.15),
        Inches(1.5),
        "唯一生成入口",
        [generation_manifest["generator"]["entry"]],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(0.75),
        Inches(3.85),
        Inches(1.5),
        "框架输入",
        [
            generation_manifest["framework_inputs"]["frontend"],
            generation_manifest["framework_inputs"]["domain"],
            generation_manifest["framework_inputs"]["backend"],
        ],
        THEME.accent_alt,
    )
    add_card(
        slide,
        Inches(5.35),
        Inches(2.6),
        Inches(7.3),
        Inches(2.4),
        "这份 PPT 重点解释",
        [
            "1. framework markdown 如何编成 framework_ir.json",
            "2. product_spec / implementation_config 如何收敛成 KnowledgeBaseProject",
            "3. contract / ui_spec / backend_spec 如何继续驱动 runtime code",
            "4. governed_symbol + governance manifest/tree 如何把代码拉回上游",
        ],
        THEME.gold,
    )

    add_outline_rect(slide, Inches(5.35), Inches(5.35), Inches(7.3), Inches(1.15), RGBColor(216, 207, 190))
    add_text(
        slide,
        Inches(5.6),
        Inches(5.66),
        Inches(6.9),
        Inches(0.45),
        "阅读方法：不要先盯某个函数，看“文件链”和“对象链”更快。",
        font_size=16,
        bold=True,
    )
    add_footer(slide, "生成依据：projects/knowledge_base_basic/generated/* 与当前 runtime / governance 主链。")


def build_overall_chain_slide(prs: Any, generation_manifest: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "01 / 总链路",
        "先看一条主链：上游真相如何逐层收敛，下游代码如何只消费编译结果",
        "这里最重要的不是文件多，而是每层职责边界不反向污染。",
    )

    chain = [
        ("Framework", "共同结构语言\nframework/*.md", THEME.accent),
        ("Product Spec", "产品真相\nproduct_spec.toml + product_spec/*.toml", THEME.accent_alt),
        ("Implementation", "实现细化\nimplementation_config.toml", THEME.gold),
        ("Compiled Project", "KnowledgeBaseProject\ncontract / spec / route", THEME.accent_dark),
        ("Code", "runtime app / pages / api", THEME.ok),
        ("Evidence", "generated/* + governance", THEME.warn),
    ]
    left = Inches(0.7)
    top = Inches(2.0)
    width = Inches(1.9)
    for index, (_, body, color) in enumerate(chain):
        x = left + Inches(2.05) * index
        add_rect(slide, x, top, width, Inches(1.0), color)
        add_text(slide, x + Inches(0.12), top + Inches(0.18), width - Inches(0.24), Inches(0.55), body, font_size=14, color=THEME.paper, bold=True)
        if index < len(chain) - 1:
            add_arrow_connector(slide, x + width, top + Inches(0.5), x + Inches(2.05), top + Inches(0.5))

    add_card(
        slide,
        Inches(0.95),
        Inches(3.55),
        Inches(5.1),
        Inches(2.2),
        "这一层层分别做什么",
        [
            "Framework：定义能力、边界、基、规则、验证。",
            "Product Spec：决定这个项目到底是哪一个具体产品。",
            "Implementation Config：决定这个产品走哪条实现路径。",
            "Compiled Project：把上游三层压成代码可以消费的稳定结构块。",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(6.35),
        Inches(3.55),
        Inches(6.05),
        Inches(2.2),
        "当前可直接核对的证据",
        [
            f"resolved_modules = {len(generation_manifest['framework_inputs']['resolved_modules'])}",
            f"configuration_effects = {len(generation_manifest['configuration_effects'])}",
            f"generated files = {len(generation_manifest['generated_files'])}",
            "所有这些都不是手写说明，而是生成产物的一部分。",
        ],
        THEME.accent_alt,
    )
    add_footer(slide, "关键判断：runtime 不是新真相源，而是 KnowledgeBaseProject 的消费者。")


def build_compiler_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "02 / 编译入口",
        "knowledge_base.py 是当前知识库模板的编译器，不是普通工具函数集合",
        "它把 framework IR、产品真相和实现细化收成一个 KnowledgeBaseProject，再继续派生 contract / spec / evidence。",
    )

    add_card(
        slide,
        Inches(0.7),
        Inches(1.75),
        Inches(3.0),
        Inches(1.75),
        "输入",
        [
            "_load_product_spec()",
            "_load_implementation_config()",
            "_resolve_framework_module()",
            "_collect_framework_closure()",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.15),
        Inches(1.75),
        Inches(4.0),
        Inches(1.75),
        "核心收敛点",
        [
            "_compile_project()",
            "-> KnowledgeBaseProject",
            "-> frontend_contract",
            "-> workbench_contract",
            "-> ui_spec / backend_spec",
        ],
        THEME.accent_dark,
    )
    add_card(
        slide,
        Inches(8.55),
        Inches(1.75),
        Inches(4.1),
        Inches(1.75),
        "证据出口",
        [
            "materialize_knowledge_base_project()",
            "-> framework_ir.json",
            "-> generation_manifest.json",
            "-> governance_manifest.json",
            "-> governance_tree.json",
        ],
        THEME.ok,
    )
    add_arrow_connector(slide, Inches(3.7), Inches(2.62), Inches(4.15), Inches(2.62))
    add_arrow_connector(slide, Inches(8.15), Inches(2.62), Inches(8.55), Inches(2.62))

    add_card(
        slide,
        Inches(1.2),
        Inches(4.15),
        Inches(11.0),
        Inches(1.85),
        "你顺着文件看时，重点盯这几个函数",
        [
            "load_knowledge_base_project()：把配置读成显式结构块。",
            "_compile_project()：把上游真相压成 KnowledgeBaseProject。",
            "_build_backend_spec() / _build_ui_spec()：把项目对象继续压成运行时规格。",
            "materialize_knowledge_base_project()：把结果固化成 generated/*。",
        ],
        THEME.gold,
    )
    add_footer(slide, "当前主入口：src/project_runtime/knowledge_base.py")


def build_effects_slide(prs: Any, generation_manifest: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "03 / 配置即功能",
        "Implementation Config 不是写着好看，它必须进入下游生效位",
        "仓库当前用 generation_manifest.configuration_effects 把这件事固化成可追踪证据。",
    )

    effects = generation_manifest["configuration_effects"]
    left_items = [
        ("frontend.renderer", effects["frontend.renderer"]["targets"][0]),
        ("frontend.style_profile", effects["frontend.style_profile"]["targets"][0]),
        ("frontend.script_profile", effects["frontend.script_profile"]["targets"][0]),
        ("backend.renderer", effects["backend.renderer"]["targets"][0]),
    ]
    right_items = [
        ("backend.transport", effects["backend.transport"]["targets"][0]),
        ("backend.retrieval_strategy", effects["backend.retrieval_strategy"]["targets"][0]),
        ("evidence.product_spec_endpoint", effects["evidence.product_spec_endpoint"]["targets"][0]),
        ("artifacts.framework_ir_json", effects["artifacts.framework_ir_json"]["targets"][0]),
    ]

    add_card(
        slide,
        Inches(0.8),
        Inches(1.75),
        Inches(3.5),
        Inches(4.5),
        "implementation_config 字段",
        [f"{name} = {effects[name]['value']}" for name, _ in left_items + right_items],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.95),
        Inches(1.75),
        Inches(3.15),
        Inches(4.5),
        "下游生效位",
        [target for _, target in left_items + right_items],
        THEME.accent_alt,
    )
    add_card(
        slide,
        Inches(8.65),
        Inches(1.75),
        Inches(3.75),
        Inches(4.5),
        "检查方式",
        [
            "build_implementation_effect_manifest(project)",
            "-> generation_manifest.configuration_effects",
            "-> validate_strict_mapping.py 逐条比较",
            "只在 bundle 里重复保存 config 本身，不算生效。",
        ],
        THEME.gold,
    )
    add_arrow_connector(slide, Inches(4.3), Inches(3.95), Inches(4.95), Inches(3.95))
    add_arrow_connector(slide, Inches(8.1), Inches(3.95), Inches(8.65), Inches(3.95))
    add_footer(slide, "关键例子：evidence.product_spec_endpoint -> backend_spec.transport.product_spec_endpoint -> runtime app route")


def build_runtime_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "04 / 运行时落点",
        "runtime code 不是再推导一次产品真相，而是只消费 contract / spec / route",
        "所以你要看代码落点时，不是从 HTML 模板往回猜，而是先看它吃哪份编译结果。",
    )

    add_card(
        slide,
        Inches(0.75),
        Inches(1.85),
        Inches(3.1),
        Inches(3.5),
        "app.py",
        [
            "build_knowledge_base_runtime_app()",
            "消费 route + backend_spec.transport",
            "挂 chat_home / list / detail / product-spec 路由",
            "高风险对象：kb.runtime.page_routes",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.15),
        Inches(1.85),
        Inches(3.9),
        Inches(3.5),
        "frontend.py",
        [
            "消费 ui_spec + frontend_contract + workbench_contract",
            "输出聊天页、知识库列表页、详情页、文档页",
            "页面结构并不自己定义 product truth",
        ],
        THEME.accent_alt,
    )
    add_card(
        slide,
        Inches(8.35),
        Inches(1.85),
        Inches(4.0),
        Inches(3.5),
        "backend.py",
        [
            "消费 backend_spec + documents + return policy",
            "暴露 library contracts / chat contracts / answer behavior",
            "高风险对象：kb.api.library_contracts / kb.api.chat_contract / kb.answer.behavior",
        ],
        THEME.ok,
    )

    add_arrow_connector(slide, Inches(3.85), Inches(3.6), Inches(4.15), Inches(3.6))
    add_arrow_connector(slide, Inches(8.05), Inches(3.6), Inches(8.35), Inches(3.6))

    add_card(
        slide,
        Inches(1.2),
        Inches(5.75),
        Inches(10.9),
        Inches(0.9),
        "关键判断",
        ["如果某个运行时函数绑定了 governed_symbol，就可以再顺着 governance manifest 把它拉回 framework / product / implementation。"],
        THEME.accent_dark,
    )
    add_footer(slide, "直接打开：src/knowledge_base_runtime/app.py, frontend.py, backend.py")


def build_reverse_governance_slide(prs: Any, governance_manifest: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "05 / 反查链",
        "代码不是靠人工记忆对应上游，而是靠 governed_symbol + governance closure",
        "这条链决定了代码改动后，仓库可以反过来检查上游有没有对应项。",
    )

    answer_behavior = _governed_object(governance_manifest, "kb.answer.behavior")
    frontend_surface = _governed_object(governance_manifest, "kb.frontend.surface_contract")

    add_card(
        slide,
        Inches(0.8),
        Inches(1.85),
        Inches(3.55),
        Inches(3.85),
        "代码打点",
        [
            "# @governed_symbol ...",
            "frontend_kernel/contracts.py",
            "knowledge_base_framework/workbench.py",
            "knowledge_base_runtime/app.py",
            "knowledge_base_runtime/backend.py",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(1.85),
        Inches(3.0),
        Inches(3.85),
        "治理闭包",
        [
            "build_governance_closure(project)",
            "-> structural_objects",
            "-> candidates",
            "-> role_bindings",
            "-> strict_zone",
            "-> upstream_closure",
        ],
        THEME.accent_dark,
    )
    add_card(
        slide,
        Inches(8.1),
        Inches(1.85),
        Inches(4.15),
        Inches(3.85),
        "治理证据",
        [
            f"kb.frontend.surface_contract: framework={len(frontend_surface['sources']['framework'])} / product={len(frontend_surface['sources']['product'])} / impl={len(frontend_surface['sources']['implementation'])}",
            f"kb.answer.behavior: framework={len(answer_behavior['sources']['framework'])} / product={len(answer_behavior['sources']['product'])} / impl={len(answer_behavior['sources']['implementation'])}",
            "这些来源会继续进入 governance_manifest.json / governance_tree.json",
        ],
        THEME.gold,
    )
    add_arrow_connector(slide, Inches(4.35), Inches(3.7), Inches(4.75), Inches(3.7))
    add_arrow_connector(slide, Inches(7.75), Inches(3.7), Inches(8.1), Inches(3.7))
    add_footer(slide, "关键文件：src/project_runtime/governance.py 与 projects/knowledge_base_basic/generated/governance_manifest.json")


def build_example_endpoint_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "06 / 例子 A",
        "拿 /api/knowledge/product-spec 这条路由，看一条完整的一一对应链",
        "这是最适合入门的例子，因为它同时穿过 route、implementation effect、backend_spec 和 runtime route。",
    )

    steps = [
        ("1", "product_spec/route.toml", "定义 route.api_prefix = /api/knowledge"),
        ("2", "implementation_config.toml", "定义 evidence.product_spec_endpoint = /api/knowledge/product-spec"),
        ("3", "knowledge_base.py", "_validate_implementation_config() 保证 endpoint stay under api_prefix"),
        ("4", "knowledge_base.py", "_build_backend_spec() 写入 backend_spec.transport.product_spec_endpoint"),
        ("5", "app.py", "@app.get(transport.product_spec_endpoint) 真正挂路由"),
        ("6", "generation_manifest.json", "configuration_effects 记录这条 effect"),
    ]

    for index, (seq, label, body) in enumerate(steps):
        x = Inches(0.65) + Inches(2.07) * index
        add_card(slide, x, Inches(2.2), Inches(1.82), Inches(2.55), f"{seq} / {label}", [body], THEME.accent if index % 2 == 0 else THEME.accent_alt)
        if index < len(steps) - 1:
            add_arrow_connector(slide, x + Inches(1.82), Inches(3.45), x + Inches(2.07), Inches(3.45))

    add_card(
        slide,
        Inches(1.2),
        Inches(5.4),
        Inches(11.0),
        Inches(0.9),
        "一句话",
        ["这里不是“代码自己决定路由”，而是“上游两层配置先决定，再由编译器校验并写进 runtime 可消费结构”。"],
        THEME.ok,
    )
    add_footer(slide, "这个例子最适合顺着读：route.toml -> implementation_config.toml -> knowledge_base.py -> app.py -> generation_manifest.json")


def build_example_answer_behavior_slide(prs: Any, governance_manifest: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "07 / 例子 B",
        "拿 kb.answer.behavior 看“从代码反查回三层上游”",
        "这是当前最能体现治理链价值的例子，因为它不是简单路由，而是回答行为本身。",
    )

    answer_behavior = _governed_object(governance_manifest, "kb.answer.behavior")
    semantic = answer_behavior["semantic"]

    add_card(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(2.9),
        Inches(3.7),
        "代码落点",
        [
            "src/knowledge_base_runtime/backend.py",
            "@governed_symbol id=kb.answer.behavior",
            "answer generation / citations / return path",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.0),
        Inches(1.9),
        Inches(2.75),
        Inches(3.7),
        "framework 来源",
        [f"{item['file']}#{item['ref_id']}" for item in answer_behavior["sources"]["framework"][:4]],
        THEME.accent_alt,
    )
    add_card(
        slide,
        Inches(7.0),
        Inches(1.9),
        Inches(2.45),
        Inches(3.7),
        "product 来源",
        [f"{item['ref_id']} <- {Path(item['file']).name}" for item in answer_behavior["sources"]["product"]],
        THEME.gold,
    )
    add_card(
        slide,
        Inches(9.7),
        Inches(1.9),
        Inches(2.65),
        Inches(3.7),
        "semantic 期望",
        [
            f"citation_style = {semantic['citation_style']}",
            f"citation_required = {semantic['citation_required']}",
            f"return_path_prefix = {semantic['return_path_prefix']}",
            f"document_path_prefix = {semantic['document_path_prefix']}",
        ],
        THEME.ok,
    )
    add_arrow_connector(slide, Inches(3.7), Inches(3.7), Inches(4.0), Inches(3.7))
    add_arrow_connector(slide, Inches(6.75), Inches(3.7), Inches(7.0), Inches(3.7))
    add_arrow_connector(slide, Inches(9.45), Inches(3.7), Inches(9.7), Inches(3.7))
    add_footer(slide, "这就是“代码对象 -> governance object -> framework/product/implementation sources” 的真实反查路径。")


def build_reading_order_slide(prs: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME.bg)
    add_section_title(
        slide,
        "08 / 怎么看",
        "如果你要自己顺着搞清楚，一般按这 3 条路径看就够了",
        "不要先盯一大坨 JSON；先确定你是在做“上游追下游”，还是“下游反查上游”。",
    )

    add_card(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(3.7),
        Inches(3.9),
        "路径一：从上往下",
        [
            "product_spec / implementation_config",
            "-> knowledge_base.py",
            "-> generation_manifest.json",
            "-> frontend.py / backend.py / app.py",
            "适合回答：这个配置最后变成了什么代码？",
        ],
        THEME.accent,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.8),
        Inches(3.7),
        Inches(3.9),
        "路径二：从下往上",
        [
            "@governed_symbol",
            "-> governance_manifest.json",
            "-> sources.framework / product / implementation",
            "适合回答：这个函数/行为到底来自哪里？",
        ],
        THEME.accent_alt,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.8),
        Inches(3.7),
        Inches(3.9),
        "路径三：查配置生效",
        [
            "implementation_config field",
            "-> generation_manifest.configuration_effects",
            "-> targets",
            "-> runtime 消费点",
            "适合回答：这是不是死配置？",
        ],
        THEME.gold,
    )
    add_footer(slide, "推荐先配合 docs/框架到代码映射与反查覆盖说明.md 一起看。")


def build_presentation() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    generation_manifest = _load_json(GENERATION_MANIFEST_PATH)
    governance_manifest = _load_json(GOVERNANCE_MANIFEST_PATH)

    prs: Any = Presentation()
    prs.slide_width = THEME.width
    prs.slide_height = THEME.height

    build_cover(prs, generation_manifest)
    build_overall_chain_slide(prs, generation_manifest)
    build_compiler_slide(prs)
    build_effects_slide(prs, generation_manifest)
    build_runtime_slide(prs)
    build_reverse_governance_slide(prs, governance_manifest)
    build_example_endpoint_slide(prs)
    build_example_answer_behavior_slide(prs, governance_manifest)
    build_reading_order_slide(prs)

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output_path = build_presentation()
    print(f"[PASS] wrote {output_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
